from langchain.tools import tool
from langgraph.graph import StateGraph, START, END
from langchain.messages import AnyMessage
from typing_extensions import TypedDict, Annotated
from typing import Dict, List
from pptx import Presentation
from pathlib import Path
import json

from .llm import model

class PPTState(TypedDict):
    topic: str
    outline: dict
    slides: List[dict]
    file_path: str


def planner_agent(state: PPTState):

    topic = state["topic"]

    SYSTEM_PROMPT = f"""
        Create a professional executive-level PowerPoint presentation outline on: {topic}

        Return the response strictly in valid JSON format.

        JSON structure:
        {{
            "presentation_title": "string",
            "slides": [
                {{
                "slide_number": 1,
                "title": "string"
                }}
            ]
        }}

        Rules:
        - Return only valid JSON
        - Do not include markdown
        - Do not include explanations
        - Do not include code blocks
        - Keep slide titles concise and professional
        - Ensure logical storytelling flow
        - Include 8-15 slides
    """

    response = model.invoke(SYSTEM_PROMPT)
    #print("PLANNER::: ", response.content)
    data = json.loads(response.content)

    return {
        "outline": data
    }

def writer_agent(state: PPTState):

    topic = state["topic"]
    outline = state["outline"]

    SYSTEM_PROMPT = f"""
        You are a professional presentation writer and slide content designer.

        You will be given a structured PPT outline (JSON format) containing:
        - presentation_title
        - slides (each with slide_number, title, and purpose)

        Your task is to convert this outline into a complete PowerPoint slide content draft.

        Topic:
        {topic}

        Outline:
        {outline}

        For each slide, generate:
        - Slide title (refined if needed but keep intent unchanged)
        - 3-6 concise bullet points (executive style)

        Rules:
        - Maintain logical flow across slides
        - Keep content crisp, professional, and presentation ready
        - Use consulting style language (clear, structured, impactful)
        - Do NOT change the overall structure or number of slides
        - Do NOT add extra slides
        - Do NOT include explanations outside slide content
        - Ensure consistency with the original outline intent
        - Do not include markdown
        - Do not include explanations
        - Do not include code blocks

        Output format (STRICT):
        Return only valid JSON in this structure:

        {{
        "slides": [
            {{
            "slide_number": 1,
            "title": "...",
            "bullets": ["...", "..."],
            }}
        ]
        }}
    """

    response = model.invoke(SYSTEM_PROMPT)
    data = json.loads(response.content)
    slides = data.get("slides", "")
    #print("WRITER2::: ", slides)

    return {
        "slides": slides
    }

def builder_agent(state: PPTState):

    outline = state["outline"]
    ppt_title = outline.get("presentation_title", "")
    slides_data = state["slides"]
    
    # Title Slide
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    #subtitle = slide.placeholders[1]

    title.text = ppt_title
    #subtitle.text = "python-pptx was here!"

    # Slides
    for slide_data in slides_data:
        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        # Title of the slide
        title_shape.text = slide_data.get("title", "")

        tf = body_shape.text_frame
        #tf.text = bullet
        for bullet in slide_data.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 2

    
    # Path

    BASE_DIR = Path(__file__).resolve().parent
    ROOT_DIR = BASE_DIR.parent
    print("Base Dir: ", BASE_DIR)
    print("Root Dir: ", ROOT_DIR.parent)

    OUTPUT_DIR = ROOT_DIR.parent / "static/ppt_files"
    OUTPUT_DIR.mkdir(exist_ok=True)

    file_path =  OUTPUT_DIR / "my_ppt.pptx"

    print("Output Dir: ", OUTPUT_DIR)
    print("File Path: ", file_path)

    prs.save(str(file_path))

    return {
        "file_path": file_path
    }


def build_graph() -> StateGraph:

    agentic_graph = StateGraph(PPTState)

    # Create Nodes
    agentic_graph.add_node("planner", planner_agent)
    agentic_graph.add_node("writer", writer_agent)
    agentic_graph.add_node("builder", builder_agent)

    # Connect Edges
    agentic_graph.set_entry_point("planner")
    agentic_graph.add_edge("planner", "writer")
    agentic_graph.add_edge("writer", "builder")
    agentic_graph.add_edge("builder", END)

    agent = agentic_graph.compile()

    return agent

ppt_agent = build_graph()